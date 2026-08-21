"""Generate Project_Plan_Kaveri_3.0_Programme_v0.4.pptx from programme plan."""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path(r"E:\MVP\Kaveri 3.0\Source Code\Kaveri 3 Plan\Project_Plan_Kaveri_3.0_Programme_v0.4.pptx")

NAVY = RGBColor(0x0B, 0x3D, 0x5C)
TEAL = RGBColor(0x1A, 0x6B, 0x7A)
ACCENT = RGBColor(0xC4, 0x5C, 0x26)
LIGHT = RGBColor(0xF0, 0xF4, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1E, 0x29, 0x33)
MUTED = RGBColor(0x5A, 0x6A, 0x78)
GREEN = RGBColor(0x1B, 0x7A, 0x4E)
ROW_ALT = RGBColor(0xE8, 0xEE, 0xF2)

TOTAL = 14


def set_run(run, text, size=18, bold=False, color=DARK, font="Calibri"):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def add_rect(slide, l, t, w, h, fill):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    return sh


def tb(slide, l, t, w, h):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    return tf


def p_add(tf, text, size=16, bold=False, color=DARK, align=PP_ALIGN.LEFT, space_after=6, first=False):
    if first:
        p = tf.paragraphs[0]
    else:
        # use first empty paragraph once
        p0 = tf.paragraphs[0]
        if not p0.runs and not (p0.text or "").strip() and not getattr(tf, "_started", False):
            p = p0
            tf._started = True
        else:
            p = tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    run = p.add_run()
    set_run(run, text, size, bold, color)
    return p


def footer_ok(slide, page, W, H, total=TOTAL):
    add_rect(slide, Inches(0), H - Inches(0.35), W, Inches(0.35), NAVY)
    tf = tb(slide, Inches(0.4), H - Inches(0.32), Inches(10), Inches(0.28))
    p_add(tf, "Kaveri 3.0 Programme Plan v0.4  |  PLAN-K3-PROG-001  |  Confidential — Steering",
          size=10, color=WHITE, space_after=0)
    tf2 = tb(slide, W - Inches(1.2), H - Inches(0.32), Inches(0.9), Inches(0.28))
    p_add(tf2, f"{page}/{total}", size=10, color=WHITE, align=PP_ALIGN.RIGHT, space_after=0)


def title_bar(slide, title, subtitle, W):
    add_rect(slide, Inches(0), Inches(0), W, Inches(0.95), NAVY)
    add_rect(slide, Inches(0), Inches(0.95), W, Inches(0.08), TEAL)
    tf = tb(slide, Inches(0.5), Inches(0.18), Inches(12), Inches(0.5))
    p_add(tf, title, size=26, bold=True, color=WHITE, space_after=0)
    if subtitle:
        tf2 = tb(slide, Inches(0.5), Inches(0.55), Inches(12), Inches(0.35))
        p_add(tf2, subtitle, size=13, color=RGBColor(0xB8, 0xD0, 0xDC), space_after=0)


def card(slide, l, t, w, h, title, lines, title_color=TEAL):
    add_rect(slide, l, t, w, h, LIGHT)
    add_rect(slide, l, t, Inches(0.08), h, title_color)
    tf = tb(slide, l + Inches(0.2), t + Inches(0.12), w - Inches(0.3), Inches(0.35))
    p_add(tf, title, size=14, bold=True, color=title_color, space_after=4)
    tf2 = tb(slide, l + Inches(0.2), t + Inches(0.45), w - Inches(0.3), h - Inches(0.55))
    for line in lines:
        p_add(tf2, line, size=12, color=DARK, space_after=3)


def add_table(slide, l, t, w, rows, col_widths=None, font_size=11):
    cols = len(rows[0])
    table_shape = slide.shapes.add_table(len(rows), cols, l, t, w, Inches(0.38 * len(rows)))
    table = table_shape.table
    if col_widths:
        for i, cw in enumerate(col_widths):
            table.columns[i].width = cw
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            run = p.add_run()
            set_run(run, str(val), size=font_size, bold=(r == 0), color=WHITE if r == 0 else DARK)
            cell.fill.solid()
            if r == 0:
                cell.fill.fore_color.rgb = NAVY
            elif r % 2 == 0:
                cell.fill.fore_color.rgb = ROW_ALT
            else:
                cell.fill.fore_color.rgb = WHITE
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    return table_shape


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    W, H = prs.slide_width, prs.slide_height
    blank = prs.slide_layouts[6]

    # 1 Title
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, W, H, NAVY)
    add_rect(s, 0, Inches(5.8), W, Inches(1.7), TEAL)
    tf = tb(s, Inches(0.8), Inches(1.8), Inches(11.5), Inches(1))
    p_add(tf, "KAVERI 3.0", size=44, bold=True, color=WHITE, space_after=8)
    tf = tb(s, Inches(0.8), Inches(2.7), Inches(11.5), Inches(1.4))
    p_add(tf, "Programme Delivery Plan", size=32, bold=True, color=RGBColor(0xD4, 0xE8, 0xEF), space_after=8)
    p_add(tf, "Phase-wise SDLC  ·  11-month Go Live  ·  Data Migration from Kaveri 1.0 / 2.0",
          size=16, color=RGBColor(0xA8, 0xC5, 0xD0), space_after=0)
    tf = tb(s, Inches(0.8), Inches(6.05), Inches(11), Inches(1.2))
    p_add(tf, "Department of Stamps & Registration  |  Government of Karnataka",
          size=16, bold=True, color=WHITE, space_after=4)
    p_add(tf, "PLAN-K3-PROG-001  ·  Version 0.4 (Draft)  ·  August 2026",
          size=13, color=RGBColor(0xD0, 0xE8, 0xE8), space_after=0)

    # 2 Agenda
    s = prs.slides.add_slide(blank)
    title_bar(s, "Agenda", "Steering review pack", W)
    footer_ok(s, 2, W, H)
    tf = tb(s, Inches(0.8), Inches(1.35), Inches(11), Inches(5.5))
    for it in [
        "1. Programme constraints & objectives",
        "2. Four delivery phases — scope",
        "3. Master timeline & Go Live dates",
        "4. Phase 1 deep-dive (21 Oct 2026)",
        "5. Phases 2–4 summary",
        "6. SDLC stages per phase",
        "7. IT Cell resources (39 posts)",
        "8. Data migration from Kaveri 1.0 / 2.0",
        "9. Governance, gates & key risks",
        "10. Decisions requested",
    ]:
        p_add(tf, it, size=18, color=DARK, space_after=10)

    # 3 Constraints
    s = prs.slides.add_slide(blank)
    title_bar(s, "Hard programme constraints", "Non-negotiables for planning", W)
    footer_ok(s, 3, W, H)
    cards = [
        (Inches(0.5), "T0 — Start", ["17 August 2026", "BR schedule kick-off", "Requirements window opens"], TEAL),
        (Inches(3.55), "Phase 1 Go Live", ["3rd week of October 2026", "Target: 21 October 2026", "~9-week all-hands sprint"], ACCENT),
        (Inches(6.6), "All modules live", ["Within 11 months", "End: 16 July 2027", "Phase 4 GL: 14 July 2027"], TEAL),
        (Inches(9.65), "Data Migration", ["Kaveri 1.0 + 2.0 in scope", "Cutover with each phase", "Exit report mandatory"], GREEN),
    ]
    for left, title, lines, color in cards:
        card(s, left, Inches(1.35), Inches(2.9), Inches(2.7), title, lines, color)
    tf = tb(s, Inches(0.5), Inches(4.4), Inches(12.3), Inches(2.5))
    p_add(tf, "Implications", size=16, bold=True, color=NAVY, space_after=8)
    for line in [
        "• Phase 1 is the pathfinder — platform (User Mgmt, eKYC, eSign, Khajane, Scanning) ships first",
        "• Phases 2–4 run in parallel after 21 Oct with squad rebalance (8 Full Stack)",
        "• Phase 3 (Document Registration + EC + migration volume) is the critical path",
        "• No phase Go Live without Migration Exit Report + legacy freeze/delta cutover",
    ]:
        p_add(tf, line, size=14, color=DARK, space_after=5)

    # 4 Objectives
    s = prs.slides.add_slide(blank)
    title_bar(s, "Business objectives", "What success looks like", W)
    footer_ok(s, 4, W, H)
    add_table(
        s, Inches(0.5), Inches(1.35), Inches(12.3),
        [
            ["#", "Objective", "Measure"],
            ["O1", "All modules live in 11 months; P1 by 3rd week Oct 2026", "P1 ≤ 21-10-2026; P4 ≤ 16-07-2027"],
            ["O2", "Statutory compliance (Registration / Stamp Acts & Rules)", "Legal / Domain template & fee lock"],
            ["O3", "Shared platform once (identity, eKYC, eSign, Khajane, scan)", "Reuse across phases — no duplicate stacks"],
            ["O4", "Deliver within IT Cell capacity", "39 posts — wave allocation"],
            ["O5", "e-Gov NFR bar", "GIGW, WCAG, STQC/security, HA/DR"],
            ["O6", "Migrate Kaveri 1.0/2.0 data safely", "Reconcile + exit report per phase"],
        ],
        [Inches(0.7), Inches(6.5), Inches(5.1)], 12,
    )

    # 5 Four phases
    s = prs.slides.add_slide(blank)
    title_bar(s, "Delivery phases — scope", "Four releases; one programme", W)
    footer_ok(s, 5, W, H)
    phases = [
        ("Phase 1", "21 Oct 2026", GREEN, [
            "Marriage (Online + Offline)", "Certified Copies (CC)", "User Management",
            "eKYC · eSign · Khajane", "Scanning · Marriage MIS/Dash", "DM-P1 masters/users/marriage",
        ]),
        ("Phase 2", "16 Mar 2027", TEAL, [
            "Stamp Duty Calculator", "Guideline Value Calculator", "Market Valuator (CVC)",
            "GIS Valuation", "E-Stamp templates", "DM-P2 rate/guideline/CVC",
        ]),
        ("Phase 3", "08 Jun 2027", NAVY, [
            "Document Registration", "Encumbrance Search (EC)", "Verify · PoA · DRO/IGRO",
            "Document MIS / Dashboard", "Filings · corrections · memo", "DM-P3 books + EC indexes",
        ]),
        ("Phase 4", "14 Jul 2027", ACCENT, [
            "Firm Registration (DRO)", "Firm MIS / Dashboard", "Firm audit views",
            "Reuse Phase 1 platform", " ", "DM-P4 firm masters",
        ]),
    ]
    for i, (name, date, color, lines) in enumerate(phases):
        left = Inches(0.35 + i * 3.2)
        add_rect(s, left, Inches(1.25), Inches(3.05), Inches(5.5), LIGHT)
        add_rect(s, left, Inches(1.25), Inches(3.05), Inches(0.95), color)
        tf = tb(s, left + Inches(0.15), Inches(1.35), Inches(2.75), Inches(0.75))
        p_add(tf, name, size=18, bold=True, color=WHITE, space_after=2)
        p_add(tf, date, size=12, color=WHITE, space_after=0)
        tf2 = tb(s, left + Inches(0.15), Inches(2.4), Inches(2.75), Inches(4.2))
        for line in lines:
            if line.strip():
                p_add(tf2, "• " + line, size=12, color=DARK, space_after=6)

    # 6 Timeline
    s = prs.slides.add_slide(blank)
    title_bar(s, "Master timeline & Go Live dates", "17 Aug 2026 → 16 Jul 2027 (11 months)", W)
    footer_ok(s, 6, W, H)
    add_table(
        s, Inches(0.4), Inches(1.25), Inches(12.5),
        [
            ["Phase", "Design", "Build", "Test", "Go Live"],
            ["1 — Marriage + Platform", "24-08 → 26-09", "25-08 → 03-10", "22-09 → 14-10", "21-10-2026"],
            ["2 — Calculators / Valuator", "22-09 → 28-11", "22-10 → 30-01", "05-01 → 28-02", "16-03-2027"],
            ["3 — Document + EC", "06-10 → 19-12", "27-10 → 24-04", "01-03 → 22-05", "08-06-2027"],
            ["4 — Firm Registration", "16-02 → 28-03", "30-03 → 13-06", "18-05 → 04-07", "14-07-2027"],
        ],
        [Inches(3.2), Inches(2.3), Inches(2.4), Inches(2.3), Inches(2.3)], 12,
    )
    tf = tb(s, Inches(0.5), Inches(3.65), Inches(12), Inches(0.35))
    p_add(tf, "Key milestones", size=14, bold=True, color=NAVY, space_after=0)
    milestones = [
        ("Aug 17", "T0 / BR start", TEAL),
        ("Oct 15", "BR window end", TEAL),
        ("Oct 21", "P1 LIVE", ACCENT),
        ("Mar 16", "P2 LIVE", TEAL),
        ("Jun 08", "P3 LIVE", TEAL),
        ("Jul 14", "P4 LIVE", GREEN),
    ]
    for i, (d, lab, color) in enumerate(milestones):
        left = Inches(0.5 + i * 2.1)
        add_rect(s, left, Inches(4.2), Inches(1.9), Inches(1.55), LIGHT)
        add_rect(s, left, Inches(4.2), Inches(1.9), Inches(0.12), color)
        tf = tb(s, left + Inches(0.12), Inches(4.4), Inches(1.7), Inches(1.2))
        p_add(tf, d, size=14, bold=True, color=NAVY, space_after=4)
        p_add(tf, lab, size=12, color=DARK, space_after=0)

    # 7 Phase 1
    s = prs.slides.add_slide(blank)
    title_bar(s, "Phase 1 deep-dive — Go Live 21 Oct 2026", "All-hands · ~9 weeks · pathfinder", W)
    footer_ok(s, 7, W, H)
    card(s, Inches(0.4), Inches(1.2), Inches(6.1), Inches(2.5), "Must at cutover", [
        "Marriage Online (eSign) + Offline (DEO)",
        "User Management (login, transfer, relieving)",
        "eKYC · eSign · Khajane payment",
        "Scanning spine · Marriage Dashboard & MIS",
        "CC MVP (after 05-Oct discussion)",
        "DM-P1 data cutover (freeze 15–20 Oct)",
    ], GREEN)
    card(s, Inches(6.7), Inches(1.2), Inches(6.1), Inches(2.5), "Build waves", [
        "W0 25-Aug–12-Sep: Platform + adapters",
        "W1 08–26 Sep: Marriage Online path",
        "W2 22-Sep–03-Oct: Offline + Scan + MIS",
        "W3 06–14 Oct: CC MVP + UAT burn-down",
        "Cutover 15–21 Oct: deploy + data + smoke",
        "P1.1 to 15-Nov: Refund/Audit/MDM finish",
    ], TEAL)
    add_table(
        s, Inches(0.4), Inches(3.95), Inches(12.5),
        [
            ["SDLC stage", "Window", "Primary owners"],
            ["BR (Marriage / User Mgmt first)", "17-08 → 15-10", "BA ×2, Domain, PO"],
            ["HLD + Tech Architecture", "24-08 → 12-09", "Architect, Tech Leads"],
            ["Development (all 8 Full Stack)", "25-08 → 03-10", "FS, Integration, DBA, UI/UX"],
            ["Testing / UAT", "22-09 → 14-10", "QA, Test Eng, Perf/Sec"],
            ["Deploy + Go Live", "15-10 → 21-10", "DevOps, Migration, Steering"],
        ],
        [Inches(4.2), Inches(3.5), Inches(4.8)], 11,
    )

    # 8 Phases 2-4
    s = prs.slides.add_slide(blank)
    title_bar(s, "Phases 2–4 — summary", "After Phase 1 capacity rebalance", W)
    footer_ok(s, 8, W, H)
    add_table(
        s, Inches(0.4), Inches(1.25), Inches(12.5),
        [
            ["Phase", "Focus", "Build start", "Go Live", "Migration"],
            ["2", "Stamp/Guideline calculators, CVC, GIS, E-Stamp", "22-10-2026", "16-03-2027", "DM-P2 rates & values"],
            ["3", "Document Registration, EC, DRO/IGRO, MIS", "27-10-2026", "08-06-2027", "DM-P3 books + EC (largest)"],
            ["4", "Firm Registration (DRO)", "30-03-2027", "14-07-2027", "DM-P4 firm data"],
        ],
        [Inches(0.9), Inches(4.8), Inches(2.0), Inches(2.0), Inches(2.8)], 12,
    )
    tf = tb(s, Inches(0.5), Inches(3.5), Inches(12.2), Inches(3.3))
    p_add(tf, "Execution notes", size=16, bold=True, color=NAVY, space_after=8)
    for line in [
        "• From 22 Oct: split squads — ~3–4 Full Stack on Phase 2; ~4–5 on Phase 3; Firm after P3 freeze",
        "• Phase 2 may use shadow compare vs As-Is calculators before hard cutover",
        "• Phase 3: pilot SRO(s) then statewide by Jun 2027; EC index pipeline starts Oct 2026",
        "• Phase 4 reuses Identity, payment, eSign, scanning, audit from Phase 1",
        "• Legal: Stamp Act 1957 + CVC Rules 2003 (P2); Registration Act 1908 + Rules 1965 + Act 47/2024 (P3)",
    ]:
        p_add(tf, line, size=14, color=DARK, space_after=6)

    # 9 SDLC
    s = prs.slides.add_slide(blank)
    title_bar(s, "SDLC stages — every phase", "Overlapping windows; gates before Go Live", W)
    footer_ok(s, 9, W, H)
    stages = [
        "1  Business Requirements",
        "2  HLD",
        "3  Technical Architecture",
        "4  SDD",
        "5  LLD",
        "6  Software Development",
        "7  Software Testing",
        "8  Software Deployment",
        "9  Go Live",
        "10 Data Migration",
    ]
    for i, name in enumerate(stages):
        left = Inches(0.3 + (i % 5) * 2.55)
        top = Inches(1.4) if i < 5 else Inches(4.15)
        color = ACCENT if i == 9 else NAVY
        add_rect(s, left, top, Inches(2.4), Inches(2.15), LIGHT)
        add_rect(s, left, top, Inches(2.4), Inches(0.5), color)
        tf = tb(s, left + Inches(0.1), top + Inches(0.7), Inches(2.2), Inches(1.2))
        p_add(tf, name, size=13, bold=True, color=NAVY, align=PP_ALIGN.CENTER, space_after=0)
        tf0 = tb(s, left + Inches(0.1), top + Inches(0.08), Inches(2.2), Inches(0.35))
        p_add(tf0, name.split()[0], size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER, space_after=0)

    # 10 Resources
    s = prs.slides.add_slide(blank)
    title_bar(s, "IT Cell resources — 39 posts", "Source: K3_ITCellRequirement.pdf", W)
    footer_ok(s, 10, W, H)
    add_table(
        s, Inches(0.5), Inches(1.2), Inches(12.3),
        [
            ["Category", "Roles", "Posts"],
            ["Leadership", "PO, PM, Transition, Content, Security, DevOps", "6"],
            ["Architecture / Tech", "Technical Architect, Technology Lead ×2", "3"],
            ["BA / Domain", "Business Analyst ×2, Domain Expert", "3"],
            ["Engineering", "Full Stack ×8, Integration, UI/UX, DBA, Migration", "12"],
            ["Quality", "QA ×2, Test Engineers ×4, Perf & Security Test Lead", "7"],
            ["Analytics", "BI & Analytics, AI/ML Specialist", "2"],
            ["Support", "L2 Support Engineers", "6"],
            ["Total", "", "39"],
        ],
        [Inches(2.5), Inches(8.3), Inches(1.5)], 12,
    )
    tf = tb(s, Inches(0.5), Inches(5.55), Inches(12.3), Inches(1.3))
    p_add(tf, "Capacity rule: All 8 Full Stack on Phase 1 until 21 Oct → then split P2 / P3 squads. Migration Specialist continuous from M1.",
          size=13, bold=True, color=NAVY, space_after=4)
    p_add(tf, "Primary bottleneck: Full Stack (8) and single Integration Engineer (Khajane → eSign → eKYC serialised).",
          size=12, color=DARK, space_after=0)

    # 11 Migration overview
    s = prs.slides.add_slide(blank)
    title_bar(s, "Data migration — Kaveri 1.0 / 2.0", "Continuous workstream §5A · MIG-K3-PROG-001", W)
    footer_ok(s, 11, W, H)
    card(s, Inches(0.4), Inches(1.25), Inches(6.1), Inches(3.0), "In scope", [
        "Structured DB data from Kaveri 1.0 & 2.0",
        "Users, offices, masters, fee/receipt keys",
        "Marriage, Document, EC indexes, Firm",
        "Already-stored scans/PDFs (re-link)",
        "2.0 = primary cutover; 1.0 = historical fill",
    ], TEAL)
    card(s, Inches(6.7), Inches(1.25), Inches(6.1), Inches(3.0), "Approach", [
        "Discover → Map → ETL → Dry-run",
        "Reconcile → UAT on migrated data",
        "Freeze + Delta → Cutover → Hypercare",
        "Idempotent pipelines + quarantine tables",
        "Rollback: legacy read-only fallback window",
    ], ACCENT)
    tf = tb(s, Inches(0.5), Inches(4.55), Inches(12.2), Inches(2.2))
    p_add(tf, "Out of scope unless Steering promotes: paper digitisation backlog (#32), Accounts (#34), Mobile (#48)",
          size=13, color=MUTED, space_after=10)
    p_add(tf, "Gate: No phase Go Live without signed Migration Exit Report (counts, samples, Domain + PO + Security).",
          size=15, bold=True, color=NAVY, space_after=0)

    # 12 Migration packages
    s = prs.slides.add_slide(blank)
    title_bar(s, "Migration packages by phase", "DM-P1 … DM-P4 must complete before respective Go Live", W)
    footer_ok(s, 12, W, H)
    add_table(
        s, Inches(0.4), Inches(1.25), Inches(12.5),
        [
            ["Phase", "Package IDs", "Must migrate", "Cutover"],
            ["1", "DM-P1-01..06", "Masters, users, marriage, images, fee refs", "15–20 Oct 2026"],
            ["2", "DM-P2-01..04", "Stamp rates, guideline, CVC, GIS keys", "With 16-03-2027 GL"],
            ["3", "DM-P3-01..07", "Regn books, pendency, EC indexes, deed links", "With 08-06-2027 GL"],
            ["4", "DM-P4-01..02", "Firm masters, certificates, open amendments", "With 14-07-2027 GL"],
        ],
        [Inches(1.0), Inches(2.2), Inches(6.3), Inches(3.0)], 12,
    )
    tf = tb(s, Inches(0.5), Inches(3.85), Inches(12.2), Inches(2.9))
    p_add(tf, "Near-term migration milestones", size=15, bold=True, color=NAVY, space_after=8)
    for line in [
        "• 05 Sep 2026 — Source inventory & access pack complete (1.0 + 2.0)",
        "• 20 Sep 2026 — Phase 1 field-mapping workbook frozen",
        "• 05 Oct 2026 — Phase 1 ETL dry-run complete; UAT only on migrated data",
        "• 20 Oct 2026 — DM-P1 Migration Exit Report signed",
        "• Oct 2026 onward — Phase 3 EC index pipeline (largest volume) starts early",
    ]:
        p_add(tf, line, size=14, color=DARK, space_after=5)

    # 13 Risks
    s = prs.slides.add_slide(blank)
    title_bar(s, "Key risks & governance", "Steering cadence: fortnightly in the 11-month window", W)
    footer_ok(s, 13, W, H)
    add_table(
        s, Inches(0.35), Inches(1.2), Inches(12.6),
        [
            ["ID", "Risk", "Mitigation"],
            ["R-00", "Phase 1 only ~9 weeks to 21 Oct", "All-hands FS; daily war-room; scope freeze"],
            ["R-00b", "P1 data migration incomplete", "DM-P1 from 01-Sep; UAT on migrated data"],
            ["R-01", "Only 8 Full Stack for 4 phases", "Squad waves; no parallel greenfield"],
            ["R-02", "Single Integration Engineer", "Stub-first; Khajane→eSign→eKYC order"],
            ["R-07", "EC/register migration volume", "EC pipeline from Oct; pilot SRO parity"],
            ["R-08", "Scope creep (mobile/paper/accounts)", "Change board; equal swap-out"],
        ],
        [Inches(1.0), Inches(4.5), Inches(7.1)], 11,
    )
    tf = tb(s, Inches(0.5), Inches(5.4), Inches(12.2), Inches(1.4))
    p_add(tf, "Go Live gate (every phase): RTM green · Legal lock · Migration Exit · Freeze/rollback · Security/Perf/DR · Training · Steering sign-off",
          size=13, bold=True, color=NAVY, space_after=0)

    # 14 Decisions
    s = prs.slides.add_slide(blank)
    title_bar(s, "Decisions requested", "Approve plan v0.4 to proceed", W)
    footer_ok(s, 14, W, H)
    tf = tb(s, Inches(0.6), Inches(1.3), Inches(12), Inches(4.8))
    for num, text in [
        ("1", "Approve Phase 1 Go Live target — 21 October 2026 (3rd week of October)"),
        ("2", "Approve 11-month programme end — all modules live by 14 July 2027 (≤ 16 July 2027)"),
        ("3", "Approve four-phase scope (Marriage platform → Calculators → Document/EC → Firm)"),
        ("4", "Approve Data Migration workstream from Kaveri 1.0/2.0 with mandatory exit reports"),
        ("5", "Confirm IT Cell 39-post roster as delivery capacity baseline"),
        ("6", "Authorise legacy write-freeze windows (Phase 1: 15–20 Oct 2026) via department circular"),
    ]:
        p_add(tf, f"{num}.  {text}", size=16, color=DARK, space_after=12)
    tf2 = tb(s, Inches(0.6), Inches(5.9), Inches(12), Inches(0.7))
    p_add(tf2, "Source: Project_Plan_Kaveri_3.0_Programme_v0.4.md / .docx", size=12, color=MUTED, space_after=0)

    # Remove broken footer helper leftovers — already using footer_ok

    prs.save(OUT)
    print(f"Saved {OUT} ({OUT.stat().st_size} bytes)")


if __name__ == "__main__":
    main()
