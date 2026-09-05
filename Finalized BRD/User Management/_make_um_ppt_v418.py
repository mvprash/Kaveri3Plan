# -*- coding: utf-8 -*-
"""Build BRD_User_Management_v4.18.pptx — a review deck for the KAVERI 3.0
User Management BRD.

Content and figures are taken from BRD_User_Management_v4.18.docx so the deck
cannot drift from the signed document: every diagram is the image embedded in
the BRD itself, extracted at build time.
"""
from __future__ import annotations

import shutil
import struct
import sys
import tempfile
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

sys.stdout.reconfigure(encoding="utf-8")

BASE = Path(__file__).resolve().parent
DOCX = BASE / "BRD_User_Management_v4.18.docx"
DST = BASE / "BRD_User_Management_v4.18.pptx"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

MARGIN = Inches(0.62)
CONTENT_W = SLIDE_W - 2 * MARGIN
BODY_TOP = Inches(1.72)
BODY_BOTTOM = Inches(6.85)

NAVY = RGBColor(0x0B, 0x25, 0x45)
BLUE = RGBColor(0x1B, 0x6C, 0xA8)
GOLD = RGBColor(0xC8, 0x9B, 0x3C)
INK = RGBColor(0x1F, 0x2A, 0x37)
MUTED = RGBColor(0x5C, 0x6B, 0x7B)
RULE = RGBColor(0xD4, 0xDD, 0xE6)
BAND = RGBColor(0xF2, 0xF6, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Segoe UI"

# BRD figure -> embedded media part (established from the document body order).
FIGURES = {
    "P-01": "image1.png",
    "S-01": "image2.png",
    "P-02": "image3.png",
    "P-03": "image4.png",
    "P-04": "image5.png",
    "P-05": "image6.png",
    "P-06": "image7.png",
    "S-02": "image8.png",
    "S-03": "image9.png",
    "S-04": "image10.png",
    "S-05": "image11.png",
    "S-06": "image12.png",
    "P-07": "image13.png",
    "P-08": "image14.png",
    "P-09": "image15.png",
    "P-12": "image16.png",
    "P-10": "image17.png",
    "P-11": "image18.png",
    "P-13": "image19.png",
}


# --------------------------------------------------------------------------
# low-level helpers
# --------------------------------------------------------------------------
def png_dimensions(data: bytes) -> tuple[int, int]:
    return struct.unpack(">II", data[16:24])


def extract_figures(target: Path) -> dict[str, Path]:
    target.mkdir(parents=True, exist_ok=True)
    paths: dict[str, Path] = {}
    with zipfile.ZipFile(DOCX) as zf:
        for tag, media in FIGURES.items():
            out = target / f"{tag}_{media}"
            out.write_bytes(zf.read(f"word/media/{media}"))
            paths[tag] = out
    return paths


def add_rect(slide, x, y, w, h, fill=None, line=None, shape=MSO_SHAPE.RECTANGLE):
    box = slide.shapes.add_shape(shape, x, y, w, h)
    if fill is None:
        box.fill.background()
    else:
        box.fill.solid()
        box.fill.fore_color.rgb = fill
    if line is None:
        box.line.fill.background()
    else:
        box.line.color.rgb = line
        box.line.width = Pt(0.75)
    box.shadow.inherit = False
    return box


def add_text(
    slide,
    x,
    y,
    w,
    h,
    text,
    size=14,
    bold=False,
    color=INK,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    italic=False,
    spacing=1.0,
):
    box = slide.shapes.add_textbox(x, y, w, h)
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    frame.vertical_anchor = anchor
    para = frame.paragraphs[0]
    para.alignment = align
    para.line_spacing = spacing
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = FONT
    return box


def add_bullets(
    slide,
    x,
    y,
    w,
    h,
    items,
    size=15,
    color=INK,
    spacing=1.18,
    space_after=9,
    marker="—",
):
    box = slide.shapes.add_textbox(x, y, w, h)
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    for index, item in enumerate(items):
        if isinstance(item, tuple):
            lead, rest = item
        else:
            lead, rest = None, item
        para = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        para.line_spacing = spacing
        para.space_after = Pt(space_after)
        if marker:
            tick = para.add_run()
            tick.text = f"{marker}  "
            tick.font.size = Pt(size)
            tick.font.bold = True
            tick.font.color.rgb = GOLD
            tick.font.name = FONT
        if lead:
            head = para.add_run()
            head.text = f"{lead}  "
            head.font.size = Pt(size)
            head.font.bold = True
            head.font.color.rgb = NAVY
            head.font.name = FONT
        run = para.add_run()
        run.text = rest
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = FONT
    return box


def set_notes(slide, text: str) -> None:
    slide.notes_slide.notes_text_frame.text = text


# --------------------------------------------------------------------------
# slide scaffolding
# --------------------------------------------------------------------------
class Deck:
    def __init__(self) -> None:
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_W
        self.prs.slide_height = SLIDE_H
        self.blank = self.prs.slide_layouts[6]
        self.number = 0

    def _bare(self):
        return self.prs.slides.add_slide(self.blank)

    def chrome(self, slide, title: str, kicker: str | None = None) -> None:
        """Standard content-slide header and footer."""
        self.number += 1
        if kicker:
            add_text(
                slide,
                MARGIN,
                Inches(0.46),
                CONTENT_W,
                Inches(0.26),
                kicker.upper(),
                size=10.5,
                bold=True,
                color=BLUE,
            )
            title_y = Inches(0.74)
        else:
            title_y = Inches(0.55)
        add_text(slide, MARGIN, title_y, CONTENT_W, Inches(0.55), title, size=27, bold=True, color=NAVY)
        add_rect(slide, MARGIN, Inches(1.46), Inches(1.15), Pt(3.2), fill=GOLD)
        add_rect(slide, MARGIN, Inches(7.02), CONTENT_W, Pt(0.9), fill=RULE)
        add_text(
            slide,
            MARGIN,
            Inches(7.12),
            Inches(9.0),
            Inches(0.26),
            "KAVERI 3.0  ·  User Management BRD v4.18  ·  Department of Stamps and Registration, Government of Karnataka",
            size=9,
            color=MUTED,
        )
        add_text(
            slide,
            SLIDE_W - MARGIN - Inches(1.0),
            Inches(7.12),
            Inches(1.0),
            Inches(0.26),
            str(self.number),
            size=9,
            bold=True,
            color=BLUE,
            align=PP_ALIGN.RIGHT,
        )

    # ---- slide types ----
    def title_slide(self) -> None:
        slide = self._bare()
        add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=NAVY)
        add_rect(slide, 0, 0, Inches(0.18), SLIDE_H, fill=GOLD)
        add_text(
            slide,
            Inches(1.15),
            Inches(1.62),
            Inches(11.0),
            Inches(0.34),
            "KAVERI 3.0  ·  DEPARTMENT OF STAMPS AND REGISTRATION, GOVERNMENT OF KARNATAKA",
            size=12,
            bold=True,
            color=GOLD,
        )
        add_text(
            slide,
            Inches(1.15),
            Inches(2.30),
            Inches(11.2),
            Inches(1.5),
            "Business Requirements Document",
            size=44,
            bold=True,
            color=WHITE,
        )
        add_text(
            slide,
            Inches(1.15),
            Inches(3.20),
            Inches(11.2),
            Inches(0.8),
            "User Management Module",
            size=32,
            color=RGBColor(0x9E, 0xC4, 0xE4),
        )
        add_rect(slide, Inches(1.15), Inches(4.20), Inches(2.1), Pt(3.2), fill=GOLD)
        add_text(
            slide,
            Inches(1.15),
            Inches(4.62),
            Inches(11.0),
            Inches(0.9),
            "Identity  ·  Passwordless authentication  ·  Sanctioned post occupancy  ·  RBAC  ·  Officer lifecycle",
            size=15,
            color=RGBColor(0xC8, 0xD7, 0xE6),
        )
        add_text(
            slide,
            Inches(1.15),
            Inches(5.85),
            Inches(11.0),
            Inches(0.9),
            "Document BRD-K3-UM-001   |   Version 4.18   |   04 September 2026\n"
            "Author: Nandha Kumar (Business Analyst)   |   Status: In review — pending Domain Expert sign-off",
            size=12.5,
            color=RGBColor(0xA9, 0xBC, 0xCE),
            spacing=1.35,
        )
        set_notes(
            slide,
            "Walkthrough of the User Management BRD v4.18 for the KAVERI 3.0 platform. "
            "All diagrams in this deck are the approved figures embedded in the BRD itself.",
        )
        self.number += 1

    def divider(self, number: str, title: str, blurb: str) -> None:
        slide = self._bare()
        self.number += 1
        add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=NAVY)
        add_rect(slide, 0, 0, Inches(0.18), SLIDE_H, fill=GOLD)
        add_text(
            slide, Inches(1.15), Inches(2.55), Inches(2.0), Inches(1.3), number, size=76, bold=True, color=RGBColor(0x1D, 0x4E, 0x80)
        )
        add_text(slide, Inches(2.55), Inches(2.86), Inches(9.6), Inches(0.9), title, size=34, bold=True, color=WHITE)
        add_rect(slide, Inches(2.55), Inches(3.86), Inches(1.6), Pt(3.2), fill=GOLD)
        add_text(
            slide, Inches(2.55), Inches(4.20), Inches(9.4), Inches(1.0), blurb, size=14.5,
            color=RGBColor(0xB6, 0xC9, 0xDA), spacing=1.28,
        )
        add_text(
            slide, SLIDE_W - MARGIN - Inches(1.0), Inches(7.02), Inches(1.0), Inches(0.26),
            str(self.number), size=9, bold=True, color=RGBColor(0x5E, 0x82, 0xA6), align=PP_ALIGN.RIGHT,
        )

    def bullets_slide(self, title, kicker, items, lead_in=None, size=15) -> None:
        slide = self._bare()
        self.chrome(slide, title, kicker)
        top = BODY_TOP
        if lead_in:
            add_text(slide, MARGIN, top, CONTENT_W, Inches(0.6), lead_in, size=14, color=MUTED, spacing=1.25)
            top = top + Inches(0.72)
        add_bullets(slide, MARGIN, top, CONTENT_W, BODY_BOTTOM - top, items, size=size)
        return slide

    def two_column_slide(self, title, kicker, left_head, left_items, right_head, right_items) -> None:
        slide = self._bare()
        self.chrome(slide, title, kicker)
        gap = Inches(0.5)
        col_w = (CONTENT_W - gap) / 2
        for index, (head, items) in enumerate(((left_head, left_items), (right_head, right_items))):
            x = MARGIN + index * (col_w + gap)
            add_rect(slide, x, BODY_TOP, col_w, Inches(0.42), fill=NAVY)
            add_text(
                slide, x + Inches(0.18), BODY_TOP + Inches(0.09), col_w - Inches(0.3), Inches(0.28),
                head, size=12.5, bold=True, color=WHITE,
            )
            add_bullets(
                slide, x, BODY_TOP + Inches(0.66), col_w, BODY_BOTTOM - BODY_TOP - Inches(0.66),
                items, size=13.5, spacing=1.16, space_after=8,
            )
        return slide

    def metrics_slide(self, title, kicker, metrics, footnote=None) -> None:
        slide = self._bare()
        self.chrome(slide, title, kicker)
        gap = Inches(0.34)
        per_row = 4
        card_w = (CONTENT_W - gap * (per_row - 1)) / per_row
        card_h = Inches(1.92)
        for index, (value, label) in enumerate(metrics):
            row, col = divmod(index, per_row)
            x = MARGIN + col * (card_w + gap)
            y = BODY_TOP + Inches(0.18) + row * (card_h + gap)
            add_rect(slide, x, y, card_w, card_h, fill=BAND, line=RULE)
            add_rect(slide, x, y, card_w, Pt(3.4), fill=GOLD)
            add_text(slide, x + Inches(0.22), y + Inches(0.34), card_w - Inches(0.44), Inches(0.7),
                     value, size=36, bold=True, color=NAVY)
            add_text(slide, x + Inches(0.22), y + Inches(1.12), card_w - Inches(0.44), Inches(0.66),
                     label, size=11.5, color=MUTED, spacing=1.15)
        if footnote:
            add_text(slide, MARGIN, Inches(6.52), CONTENT_W, Inches(0.4), footnote, size=11.5,
                     color=MUTED, italic=True)
        return slide

    def table_slide(self, title, kicker, headers, rows, col_widths, lead_in=None,
                    header_size=11.5, body_size=11, row_h=Inches(0.38)) -> None:
        slide = self._bare()
        self.chrome(slide, title, kicker)
        top = BODY_TOP
        if lead_in:
            add_text(slide, MARGIN, top, CONTENT_W, Inches(0.5), lead_in, size=13.5, color=MUTED, spacing=1.22)
            top = top + Inches(0.60)

        total = sum(col_widths)
        widths = [Inches(13.333 - 1.24) * (w / total) for w in col_widths]
        shape = slide.shapes.add_table(len(rows) + 1, len(headers), MARGIN, top,
                                       CONTENT_W, row_h * (len(rows) + 1))
        table = shape.table
        table.first_row = True
        table.horz_banding = False
        for index, width in enumerate(widths):
            table.columns[index].width = Emu(int(width))

        for index, head in enumerate(headers):
            cell = table.cell(0, index)
            cell.text = head
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY
            cell.margin_left = Inches(0.1)
            cell.margin_right = Inches(0.1)
            cell.margin_top = Inches(0.05)
            cell.margin_bottom = Inches(0.05)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(header_size)
            para.font.bold = True
            para.font.color.rgb = WHITE
            para.font.name = FONT

        for r, row in enumerate(rows, start=1):
            for c, value in enumerate(row):
                cell = table.cell(r, c)
                cell.text = str(value)
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r % 2 else BAND
                cell.margin_left = Inches(0.1)
                cell.margin_right = Inches(0.1)
                cell.margin_top = Inches(0.04)
                cell.margin_bottom = Inches(0.04)
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                para = cell.text_frame.paragraphs[0]
                para.font.size = Pt(body_size)
                para.font.color.rgb = INK
                para.font.name = FONT
                para.font.bold = c == 0 and len(headers) > 2
        return slide

    def figure_slide(self, tag, name, title, kicker, refs, figures, takeaways=None, notes=None) -> None:
        slide = self._bare()
        self.chrome(slide, title, kicker)
        add_text(slide, MARGIN, Inches(1.62), CONTENT_W, Inches(0.3),
                 f"{tag} {name}   ·   {refs}", size=11.5, bold=True, color=BLUE)

        takeaways = takeaways or []
        strip_h = Inches(0.34) * len(takeaways) + (Inches(0.16) if takeaways else Inches(0))
        frame_top = Inches(2.06)
        frame_h = BODY_BOTTOM - frame_top - strip_h
        frame_w = CONTENT_W

        path = figures[tag]
        px_w, px_h = png_dimensions(path.read_bytes())
        scale = min(frame_w / px_w, frame_h / px_h)
        draw_w, draw_h = int(px_w * scale), int(px_h * scale)
        x = MARGIN + (frame_w - draw_w) // 2
        y = frame_top + (frame_h - draw_h) // 2

        add_rect(slide, x - Inches(0.06), y - Inches(0.06),
                 draw_w + Inches(0.12), draw_h + Inches(0.12), fill=WHITE, line=RULE)
        slide.shapes.add_picture(str(path), x, y, draw_w, draw_h)

        if takeaways:
            add_bullets(slide, MARGIN, BODY_BOTTOM - strip_h + Inches(0.10), CONTENT_W,
                        strip_h, takeaways, size=12, spacing=1.05, space_after=3)
        if notes:
            set_notes(slide, notes)
        return slide


# --------------------------------------------------------------------------
# deck content
# --------------------------------------------------------------------------
def build(figures: dict[str, Path]) -> Presentation:
    deck = Deck()

    deck.title_slide()

    deck.bullets_slide(
        "Agenda",
        "Contents",
        [
            ("1.", "Introduction, scope and business objectives"),
            ("2.", "Identity and passwordless authentication — Citizens, DSR Officers, Other Department users"),
            ("3.", "Posts, roles and role-based access control"),
            ("4.", "DSR Officer lifecycle — creation, transfer, occupancy and temporary absence"),
            ("5.", "Non-functional requirements, reporting, acceptance and sign-off"),
        ],
        lead_in="A walkthrough of BRD-K3-UM-001 v4.18. Every process diagram shown is the approved figure "
                "embedded in the BRD; the editable .drawio sources sit under ProcessDiagrams/User_Management.",
        size=16,
    )

    deck.table_slide(
        "Document control",
        "BRD-K3-UM-001",
        ["Field", "Value"],
        [
            ["Document ID", "BRD-K3-UM-001"],
            ["Version", "4.18  (04 September 2026)"],
            ["Status", "In review — pending Domain Expert sign-off"],
            ["Module", "User Management"],
            ["Author (BA)", "Nandha Kumar"],
            ["Product Owner", "Prashanth"],
            ["Domain expert / reviewer", "Prabhakar Naik"],
            ["Target audience", "Kaveri IT Cell, Department of Stamps and Registration, Government of Karnataka"],
            ["Legal basis (primary)", "Information Technology Act 2000; Indian Registration Act 1908; Aadhaar Act 2016"],
            ["State rules (primary)", "Karnataka e-Governance hosting and security norms; MeitY / CERT-In / STQC / GIGW"],
        ],
        col_widths=[3.0, 9.1],
        row_h=Inches(0.44),
    )

    # ---------------- Section 1 ----------------
    deck.divider(
        "01",
        "Introduction and Scope",
        "Why the module exists, what it replaces, and where its boundaries lie.",
    )

    deck.two_column_slide(
        "Purpose and background",
        "Sections 1.1 – 1.2",
        "PURPOSE",
        [
            "Defines the business requirements for the KAVERI 3.0 User Management module.",
            "Covers identity, passwordless authentication, sanctioned post occupancy, RBAC and officer lifecycle.",
            "Applies to Citizens, DSR Officers and Other Department users.",
            "Is the agreed basis for design, development, testing and sign-off.",
        ],
        "BACKGROUND",
        [
            "KAVERI 3.0 needs one platform service for users, roles, posts and module access.",
            "Replaces fragmented Kaveri 2.0 user administration.",
            "One User Master and one Role Master — no separate stores per category.",
            "OTP-only login with no passwords; post-based DSR access; Application Admin maintains privilege mapping.",
        ],
    )

    deck.bullets_slide(
        "Scope",
        "Section 1.3",
        [
            "User registration and profile management for Citizens, DSR Officers and Other Department users (single User Master).",
            "Passwordless authentication, session policy, Citizen lost-mobile reset, DSR post selection and additional charge (§6.2–6.5).",
            "Unified Role Master; Posts and Sanctioned Posts masters; office and officer hierarchies; RBAC via Module, Function and Resource mapping (§6.5).",
            "DSR lifecycle: post assignment, Transfer Out / In, occupancy refresh, Temporary Absence and Temporary Charge (§6.6).",
            "Administrative user management, audit logging and reporting (§6.7–8).",
        ],
        lead_in="In scope for the User Management module:",
        size=15,
    )

    deck.bullets_slide(
        "Business objectives",
        "Section 3",
        [
            "Provide a secure and reliable mechanism for users to register and access the system.",
            "Enable administrators to manage user accounts and access rights efficiently.",
            "Enforce role-based access control to protect sensitive data and functionality.",
            "Maintain the DSR Officer Hierarchy Master aligned to the departmental organisation chart.",
            "Support Module Function and Resource masters with Role–Module–Function mapping and runtime enforcement.",
            "Reduce support overhead on account access issues through reliable OTP delivery.",
            "Ensure compliance with organisational security and data privacy policies.",
            "Provide auditability of user actions for security and compliance reporting.",
        ],
        size=14.5,
    )

    deck.metrics_slide(
        "The module at a glance",
        "Scale of the specification",
        [
            ("84", "Functional requirements, FR-UM-001 to FR-UM-084"),
            ("3", "User categories in a single User Master"),
            ("19", "Approved process and structure diagrams"),
            ("8", "Business modules under the Module Master"),
            ("40", "Posts in the Posts Master across 8 divisions"),
            ("0", "Passwords stored — OTP-only authentication"),
            ("24", "Non-functional requirements across 6 categories"),
            ("7 yrs", "Minimum audit-log retention"),
        ],
        footnote="Counts taken from BRD v4.18. Detailed requirement text remains authoritative in Section 6 of the document.",
    )

    deck.table_slide(
        "Stakeholders",
        "Section 4",
        ["Name / Role", "Department", "Responsibility"],
        [
            ["Kaveri IT Cell", "Engineering", "Reviews technical feasibility"],
            ["Citizens (Public users)", "External", "Self-register and access citizen portal services"],
            ["DSR Officers & Other Department users", "Government", "Access departmental modules via OTP + biometrics"],
            ["Prashanth", "Product Owner", "Owns scope and signs off UAT"],
            ["Prabhakar Naik", "Domain Expert", "Confirms departmental rules, hierarchies and role mapping"],
        ],
        col_widths=[3.4, 2.6, 6.1],
        row_h=Inches(0.5),
    )

    # ---------------- Section 2 ----------------
    deck.divider(
        "02",
        "Identity and Authentication",
        "One User Master, three categories, and no passwords anywhere in the platform.",
    )

    deck.table_slide(
        "Three user categories, one User Master",
        "Section 6.5.2",
        ["User category", "Username", "Authentication", "Lost-mobile reset"],
        [
            [
                "Public (Citizen)",
                "Preferred Username chosen at registration (FR-UM-062)",
                "Username + Captcha + OTP to mobile",
                "Self-service: 3 of 5 security questions + PIN to registered email + OTP to new mobile (FR-UM-056)",
            ],
            [
                "DSR Officer",
                "KGID (FR-UM-062, FR-UM-064)",
                "Username (KGID) + Captcha + OTP to mobile + Biometrics, then post selection (FR-UM-052)",
                "Not available — administrator changes the mobile with reason and audit (FR-UM-065)",
            ],
            [
                "Other Department",
                "KGID (FR-UM-062, FR-UM-064)",
                "Username (KGID) + Captcha + OTP to mobile + Biometrics",
                "Not available — administrator changes the mobile with reason and audit (FR-UM-065)",
            ],
        ],
        lead_in="Differentiation is by User Category on users and Role Category on roles — there are no separate masters per category.",
        col_widths=[1.9, 2.7, 3.6, 4.0],
        body_size=10.5,
        row_h=Inches(1.0),
    )

    deck.figure_slide(
        "S-01", "Access Model",
        "How identity resolves to access", "Section 6.5.1",
        "FR-UM-004, FR-UM-009, FR-UM-029, FR-UM-034, FR-UM-052, FR-UM-062",
        figures,
        takeaways=[
            "Citizens hold a Citizen role; Other Department users hold exactly one role; DSR Officers derive roles from post occupancy via Post–Role mapping.",
            "The Username is the single unique login identifier across the whole User Master — email and mobile carry no uniqueness constraint.",
        ],
        notes="One User Master and one Role Master. Category drives the creation path, the authentication factors and the eligible roles.",
    )

    deck.figure_slide(
        "P-01", "Citizen Self-Registration",
        "Citizen self-registration", "Section 6.1.1",
        "FR-UM-001, FR-UM-055, FR-UM-062, FR-UM-063",
        figures,
        takeaways=[
            "Instant self-registration with no approval workflow; the citizen chooses their own Username, checked for availability across the entire User Master.",
            "Email and mobile are each verified by a separate OTP before the account is created; five security questions are captured for lost-mobile reset only.",
        ],
        notes="Steps: open registration and enter particulars; choose Username and check availability (FR-UM-062); enter email and mobile; "
              "two separate OTPs dispatched (FR-UM-063); enter both OTPs; select five security questions and answers (FR-UM-055); "
              "account created with Citizen role and answers stored hashed.",
    )

    deck.figure_slide(
        "P-02", "Login — All User Categories",
        "Login across all categories", "Section 6.5.2.1",
        "FR-UM-004–FR-UM-007, FR-UM-010, FR-UM-011, FR-UM-062",
        figures,
        takeaways=[
            "The account resolves from the Username alone, so no category selection is needed at the login screen.",
            "The login OTP is never emailed — it is sent only by SMS to the registered mobile (FR-UM-010).",
        ],
        notes="Captcha is validated before OTP dispatch for all categories (FR-UM-011). DSR and Other Department users add biometrics; "
              "DSR Officers then continue to post selection (FR-UM-052). A Citizen who cannot receive the OTP leaves login for the lost-mobile reset flow.",
    )

    deck.figure_slide(
        "P-03", "Citizen Lost-Mobile Reset",
        "Citizen lost-mobile reset", "Section 6.5.2.2",
        "FR-UM-056 — three of five questions plus a PIN to the registered email",
        figures,
        takeaways=[
            "Two independent proofs are required before the mobile may be changed: three randomly selected security questions and a single-use PIN sent to the registered email.",
            "The new mobile number is itself OTP-verified before the change takes effect; the completed change is notified to the email and audit-logged.",
        ],
        notes="Available only to Public users (Citizens). Failed answers and PIN entries are rate-limited; five failures lock the reset flow for thirty minutes.",
    )

    deck.figure_slide(
        "P-04", "Departmental Mobile Change",
        "Departmental mobile change — administrator only", "Section 6.5.2.3",
        "FR-UM-065 — no self-service path for departmental users",
        figures,
        takeaways=[
            "DSR Officers and Other Department users hold no security questions and are never offered the FR-UM-056 reset path.",
            "An authorised administrator records a mandatory reason; the new number is OTP-verified and the change is audit-logged and notified on official email.",
        ],
        notes="This keeps departmental authentication factors under departmental control and removes any pre-login route to redirect an officer's OTP.",
    )

    deck.table_slide(
        "OTP and session policy",
        "FR-UM-069 – FR-UM-076",
        ["Control", "Rule", "Requirement"],
        [
            ["Login OTP validity", "5 minutes from dispatch (IST)", "FR-UM-069"],
            ["Registration OTP / reset PIN validity", "10 minutes from dispatch", "FR-UM-069"],
            ["Code length", "6 numeric digits", "FR-UM-070"],
            ["Incorrect entries per code", "Maximum 3, then the code is invalidated", "FR-UM-071"],
            ["Resend cooldown", "30 seconds; maximum 3 resends per channel per 15 minutes", "FR-UM-072"],
            ["Failed-login lockout", "5 failures lock the Username for 15 minutes", "FR-UM-073"],
            ["Idle timeout", "15 minutes of inactivity ends the session", "FR-UM-074"],
            ["Absolute session limit", "8 hours from login, even if the user is active", "FR-UM-075"],
            ["Concurrent sessions", "One per Username — last login wins", "FR-UM-076"],
        ],
        lead_in="Applies to every user category. The 15-minute idle timeout matters most on shared SRO counter machines, where the session post is fixed for the whole session.",
        col_widths=[3.5, 6.6, 2.0],
        row_h=Inches(0.4),
    )

    # ---------------- Section 3 ----------------
    deck.divider(
        "03",
        "Posts, Roles and Access Control",
        "Access follows the post an officer occupies — not a privilege attached to the person.",
    )

    deck.figure_slide(
        "P-05", "DSR Login Post Selection",
        "DSR Officer post selection at login", "Section 6.5.2.4",
        "FR-UM-052 — mandatory when more than one occupancy is active",
        figures,
        takeaways=[
            "One active occupancy auto-selects and skips the screen; two or more force a mandatory choice labelled Role — Post Name — Office Name (Office Code).",
            "Claims derive from the selected post only via Post–Role mapping — never from the union of other posts the officer holds.",
        ],
        notes="Relieved occupancies and Transfer In occupancies before their Joining Date are excluded from the list "
              "(FR-UM-058, FR-UM-061, FR-UM-068). Additional charge is not offered during the login flow.",
    )

    deck.figure_slide(
        "P-06", "Additional Charge After Login",
        "Additional charge of an unoccupied subordinate post", "Section 6.5.2.5",
        "FR-UM-053, FR-UM-054 — post-login, same office, no order required",
        figures,
        takeaways=[
            "A subordinate post qualifies only when it is wholly unoccupied at that office (Occupied = 0); any occupant blocks both the post and the cascade beneath it.",
            "While additional charge is active, privileges come from that post alone — a Sub-Registrar acting as FDA loses SR signing until switch-back.",
        ],
        notes="Temporary in-office cover for the session, taken after login without logout and without a Transfer Order. "
              "At most one additional charge post at a time; selection, switch-back and reversion are all audit-logged.",
    )

    deck.figure_slide(
        "S-02", "Two Vacancy Tests",
        "Two distinct vacancy tests", "Section 6.5.3",
        "FR-UM-066(a) available capacity  vs  FR-UM-066(b) wholly unoccupied",
        figures,
        takeaways=[
            "Available capacity (Occupied < Sanctioned Strength) governs post assignment at user creation and Transfer In.",
            "Wholly unoccupied (Occupied = 0) governs additional charge only — a post with strength 2 and one occupant has capacity but is not unoccupied.",
        ],
        notes="Occupied includes active occupancies and reserved Transfer In occupancies (FR-UM-067). An occupancy with a recorded "
              "Relieving Date still counts as occupied until the midnight refresh job de-allocates it.",
    )

    deck.table_slide(
        "Role Master — DSR roles by division",
        "Sections 6.5.3 – 6.5.4",
        ["Division (FR-UM-077)", "Roles under Role Category = DSR"],
        [
            ["Secretariat", "ACS / Principal Secretary / Secretary"],
            ["Top Management", "IGR"],
            ["Admin, Law & Computers", "DIGR (Admin, Law & Computers), AIGR (Admin), HQA (Admin), Sub Registrar (Admin), FDA (Admin), SDA (Admin), Typist (Admin), HQA (RTI), FDA (RTI), SDA (RTI), Statistical Inspector, Accountant Superintendent (Admin)"],
            ["Vigilance", "DIGR (Vigilance), Law Officer"],
            ["Computers", "AIGR (Computers), System Integrator, PMU, Application Developer, HQA / Project Manager (Comp), Sub Registrar (Computers), FDA (Computers), SDA (Computers)"],
            ["Enforcement", "DIGR (Enforcement), DRO, HQA (Enforcement), Sub-Registrar (SR), FDA (Enforcement), SDA (Enforcement), DEO"],
            ["Intelligence & Audit", "DIGR (Intelligence), AIGR (Audit), HQA (Audit), Superintendent (Audit), FDA (Audit), SDA (Audit), Typist (Audit)"],
            ["DIGR CVC", "DIGR CVC, JD Town Planning"],
        ],
        lead_in="Divisions come from the Division Master, not free text. Role names must be unique — FDA (Admin) and FDA (Enforcement) are distinct roles, never a shared FDA.",
        col_widths=[2.6, 9.5],
        body_size=10,
        row_h=Inches(0.40),
    )

    deck.table_slide(
        "The privilege chain",
        "Section 6.5.6",
        ["Level", "Master", "Maintained by", "Purpose"],
        [
            ["1", "Module Master", "Application Admin", "Business modules — Registration of Documents, Marriage Registration, Encumbrance Search, Certified Copy, Stamp Duty / Payments, Firm / Society Registration, User Management, MIS / Dashboards"],
            ["2", "Module Function Master", "Application Admin", "Functions under each module — VIEW, ADD, EDIT, APPROVE, SIGN, PRINT, APPLY, ISSUE, ADMIN"],
            ["3", "Resource Master", "Application Admin", "APIs and URLs linked to each Module Function, with an Is Public flag for unauthenticated endpoints"],
            ["4", "Role–Module–Function mapping", "Application Admin", "Which roles may perform which Module Functions; role names must match the Role Master exactly (FR-UM-050)"],
        ],
        lead_in="Access is modelled as User → Role(s) → Module Function(s) → Resource(s). DSR organisational roles are never named after application services.",
        col_widths=[0.8, 2.6, 2.2, 6.5],
        body_size=10.5,
        row_h=Inches(0.86),
    )

    deck.figure_slide(
        "S-03", "Runtime Enforcement",
        "Runtime access enforcement", "Section 6.5.6",
        "FR-UM-038, FR-UM-041, FR-UM-050, FR-UM-051",
        figures,
        takeaways=[
            "Every API or URL request is matched against the Resource Master and the session's Role–Module–Function claims.",
            "Deny by default — an unmatched request returns HTTP 403 and the attempt is audited. Only resources flagged Is Public bypass authentication.",
        ],
        notes="Application Admin is a system-level actor outside the Role Master and Role–Module–Function mapping (FR-UM-051), "
              "and maintains all masters in this section. Maker-checker is not required.",
    )

    deck.figure_slide(
        "S-04", "Officer Hierarchy Tree",
        "DSR Officer Hierarchy Master", "Section 6.5.7",
        "FR-UM-043 — each post's immediate parent drives relieving and Transfer In authority",
        figures,
        takeaways=[
            "Authority is granted by immediate parentage only: IGR may act on DIGR, DIGR (Enforcement) on DRO, DRO on Sub-Registrar, SR on FDA and DEO.",
            "Seeing a post lower in the tree does not confer authority over it — IGR cannot relieve a Sub-Registrar directly.",
        ],
    )

    deck.figure_slide(
        "S-05", "Office Hierarchy Tree",
        "Office Hierarchy Master", "Section 6.5.8",
        "FR-UM-059 — MS Building → IGR Head Office → District Registrar Offices → Sub-Registrar Offices",
        figures,
        takeaways=[
            "Maintained separately from the Officer Hierarchy; it scopes which offices a superior can see for Transfer Out and Transfer In.",
        ],
    )

    deck.figure_slide(
        "S-06", "Span vs Action",
        "Office span is visibility, not authority", "Section 6.5.8",
        "FR-UM-057, FR-UM-059, FR-UM-043",
        figures,
        takeaways=[
            "Both tests must pass: the target office must fall in the actor's office span, and the actor's session Post must be the immediate parent of the target Post.",
        ],
        notes="Grey shading marks offices visible in the tree; a bold border marks posts the actor may actually relieve or Transfer In. "
              "Seeing a descendant office does not grant authority over the posts in it.",
    )

    # ---------------- Section 4 ----------------
    deck.divider(
        "04",
        "Officer Lifecycle",
        "Creation, transfer, the midnight occupancy refresh, and temporary absence.",
    )

    deck.figure_slide(
        "P-07", "DSR Officer User Creation",
        "DSR Officer creation with post assignment", "Section 6.6.1",
        "FR-UM-017, FR-UM-030, FR-UM-045–048, FR-UM-062, FR-UM-066(a)",
        figures,
        takeaways=[
            "KGID is entered first because it becomes the Username; at least one sanctioned post with available capacity must be assigned or the save is blocked.",
            "Roles are shown through Post–Role mapping — there is no Primary/Secondary model, and no security questions are captured for this category.",
        ],
        notes="An optional End Date with a Deputation Reason covers time-bound deputations only; it is not a substitute for Transfer Out relieving (FR-UM-057). "
              "Biometrics capture is mandatory.",
    )

    deck.figure_slide(
        "P-08", "Other Department User Creation",
        "Other Department user creation", "Section 6.6.2",
        "FR-UM-029, FR-UM-033, FR-UM-034, FR-UM-062, FR-UM-064",
        figures,
        takeaways=[
            "Exactly one role is selected from the Role Master filtered to Role Category = Other Department; no role means no save.",
            "An optional Account End Date deactivates the user and blocks login on that date, with an audit entry.",
        ],
        notes="Stored in the same User Master with User Category = Other Department. Parent department and designation are captured; "
              "biometrics are mandatory; an authorisation letter or NOC is recommended but not mandatory to save.",
    )

    deck.figure_slide(
        "P-09", "Transfer Out / Relieving",
        "Transfer Out and relieving", "Section 6.6.3",
        "FR-UM-057, FR-UM-058, FR-UM-059, FR-UM-043, FR-UM-067, FR-UM-068",
        figures,
        takeaways=[
            "The superior sees only offices in their office span, and within those only occupancies where their session Post is the immediate parent.",
            "The occupancy stays active until 23:59 IST of the Relieving Date; the midnight job then de-allocates it and adjusts the occupied count.",
        ],
        notes="Relieving Date and Relieving Order are mandatory. If a reserved Transfer In activates in the same job run the occupied count is "
              "unchanged; otherwise it decreases by one. An officer left with no occupancies has login blocked or limited per policy.",
    )

    deck.figure_slide(
        "P-10", "Transfer In",
        "Transfer In", "Section 6.6.4",
        "FR-UM-060, FR-UM-061, FR-UM-066(a), FR-UM-067",
        figures,
        takeaways=[
            "A post with available capacity is selectable without any prior relieving; a post at full strength may receive a future-dated Transfer In only when relieving is already recorded.",
            "Recording it reserves capacity immediately, so a second Transfer In cannot take the same slot.",
        ],
        notes="Transfer / Reporting Order and Joining Date are mandatory. The incoming officer cannot log in under that post before "
              "12:00 AM IST on the Joining Date (FR-UM-061); the midnight job activates the reserved occupancy that day.",
    )

    deck.figure_slide(
        "P-12", "Handover Timeline",
        "Worked example — handover with reservation", "Section 6.6.4",
        "DRO Bengaluru → SRO Yeshwanthapura, 31-Aug to 01-Sep-2026",
        figures,
        takeaways=[
            "Relieving on 31-Aug and Transfer In joining 01-Sep coexist: capacity is reserved on the day the Transfer In is recorded, and the occupied count stays at 1 throughout.",
            "IGR at Head Office cannot perform this Transfer In, because Sub-Registrar does not report immediately to IGR.",
        ],
    )

    deck.figure_slide(
        "P-11", "Occupancy Refresh Job",
        "The midnight occupancy refresh job", "Section 6.6.5",
        "FR-UM-068 — idempotent, runs shortly after 12:00 AM IST",
        figures,
        takeaways=[
            "Occupancy changes must not wait for a user to log in: relieving, End Dates and reserved Transfer In activations all take effect on the correct calendar day.",
            "The job recalculates occupied counts and the wholly-unoccupied flag, refreshes login post selection, writes an audit record, and alerts on failure.",
        ],
    )

    deck.figure_slide(
        "P-13", "Temporary Absence and Temporary Charge",
        "Temporary absence — Leave, OOD and cover", "Section 6.6.6",
        "FR-UM-079 – FR-UM-084",
        figures,
        takeaways=[
            "Absence is recorded by the hierarchy superior only and does not free the slot — Occupied is unchanged and Transfer In must not treat the post as vacant.",
            "While an effective absence exists the officer's login is blocked entirely; the superior may assign temporary charge to a peer under them, including at another office.",
        ],
        notes="Example: a District Registrar records Leave for the Sub-Registrar of SRO Yeshwanthapura and gives temporary charge to the "
              "Sub-Registrar of SRO Jayanagar. The cover officer sees a Temporary charge row at post selection and receives the covered post's "
              "claims when it is chosen. Distinct from FR-UM-053 additional charge, which requires Occupied = 0. "
              "Full HRMS leave balances and payroll are out of scope.",
    )

    # ---------------- Section 5 ----------------
    deck.divider(
        "05",
        "Quality, Reporting and Sign-off",
        "Non-functional requirements, reports, acceptance criteria and open risks.",
    )

    deck.two_column_slide(
        "Non-functional requirements",
        "Section 7 — 24 requirements across 6 categories",
        "SECURITY AND PERFORMANCE",
        [
            "No password storage anywhere; the login OTP goes only to the registered mobile.",
            "TLS 1.2 or higher for all data in transit; biometric handling complies with the Aadhaar Act 2016 and UIDAI guidelines.",
            "Security-question answers stored hashed, never displayed, never entered by an administrator.",
            "OTP dispatch within 5 seconds; authentication completes within 2 seconds after verification.",
            "Minimum 500 concurrent authenticated sessions at launch, validated in performance testing.",
        ],
        "AVAILABILITY, AUDIT AND COMPLIANCE",
        [
            "The occupancy refresh job completes before the first login of the day; failure raises an operational alert.",
            "Disaster recovery targets RPO ≤ 24 hours and RTO ≤ 4 hours.",
            "All create, update, delete, login, recovery, access-control and transfer actions logged with timestamp and actor.",
            "Audit logs retained for at least 7 years unless a longer DSR records-retention rule applies.",
            "Bilingual Kannada and English citizen UI, GIGW accessibility, data residency in India, DPDP Act 2023 compliance.",
        ],
    )

    deck.bullets_slide(
        "Reporting requirements",
        "Section 8",
        [
            "Active, inactive and suspended users; role and permission assignments across all users.",
            "Audit log of login attempts, successful and failed, over a selected date range.",
            "Sanctioned post occupancy — sanctioned strength, occupied count, remaining capacity and wholly-unoccupied flag per office (FR-UM-066).",
            "Role-to-Module mapping showing which modules and functions each role holds.",
            "A single contact-change and recovery report covering Citizen lost-mobile recovery, administrator-initiated mobile changes and email changes.",
            "Additional charge report (FR-UM-053) and Temporary Absence / Temporary Charge report (FR-UM-079–084).",
            "Occupancy-refresh report per midnight run, with before and after occupied counts (FR-UM-068).",
            "Transfer Out / Transfer In history and officer posting / service history.",
        ],
        size=14,
    )

    deck.bullets_slide(
        "Acceptance criteria",
        "Section 9",
        [
            ("FR-UM-052", "an officer with two occupancies must choose a post before home; a single occupancy auto-selects."),
            ("FR-UM-053", "additional charge lists only wholly unoccupied posts at the same office; privileges switch until switch-back."),
            ("FR-UM-058 / 068", "relieving holds through 23:59 IST of the Relieving Date; the midnight job then de-allocates and updates counts."),
            ("FR-UM-061 / 067", "future-dated Transfer In against a full post requires recorded relieving; login is blocked until the Joining Date."),
            ("FR-UM-069–076", "OTP validity, code length, lockout, idle timeout, 8-hour session cap and single active session all enforced."),
            ("FR-UM-079–084", "absence blocks login, keeps Occupied unchanged, and temporary charge appears for the cover officer until the to_date."),
            ("Role mapping", "Role–Module–Function mapping completed and Domain Expert confirmed for every active role and module before go-live."),
        ],
        lead_in="Accepted when Section 6 is implemented and passes QA including these UAT scenarios, Section 7 is verified, "
                "the Product Owner signs off UAT, and security review closes with no critical or high findings.",
        size=13.5,
    )

    deck.table_slide(
        "Key risks and mitigations",
        "Section 10 — 24 risks logged",
        ["Risk", "Impact", "Mitigation"],
        [
            ["Role–Module–Function mapping incomplete at go-live", "High", "Only a handful of roles have seeded mappings; Application Admin with Domain Expert sign-off must complete the matrix for every active role and module."],
            ["Unauthorised relieving or Transfer In", "High", "Scope to offices under the actor (FR-UM-059), then immediate-parent posts only (FR-UM-043, FR-UM-057)."],
            ["Vacancy tests confused in build", "High", "FR-UM-066 defines available capacity and wholly unoccupied separately; additional charge uses Occupied = 0."],
            ["Double-booking a full post during handover", "High", "FR-UM-067 reserves capacity when relieving is recorded, so a second Transfer In cannot take the slot."],
            ["Leave / OOD treated as a vacancy", "High", "FR-UM-081 keeps Occupied; temporary charge is superior-assigned and distinct from FR-UM-053."],
            ["Account takeover via lost-mobile reset", "High", "Two independent proofs, OTP verification of the new number, rate limiting, email notification and audit."],
            ["Re-authentication may slow high-volume SRO counters", "Medium", "Measure end-to-end re-authentication time in performance testing with Kaveri IT Cell."],
        ],
        col_widths=[4.0, 1.0, 7.1],
        body_size=10,
        row_h=Inches(0.62),
    )

    slide = deck.bullets_slide(
        "Approval and next steps",
        "Section 12",
        [
            "Domain Expert review of the officer hierarchy, office hierarchy and Post–Role mapping seed data.",
            "Application Admin to complete Role–Module–Function mapping for every active role and module before go-live.",
            "Department to confirm the PII retention and purge policy for dormant Citizen accounts.",
            "Performance testing with Kaveri IT Cell to validate peak-load and re-authentication timings at SRO counters.",
            "Product Owner sign-off on UAT, then baseline the BRD for design and development.",
        ],
        lead_in="Signatories: Prashanth (Product Owner) · Prabhakar Naik (Domain Expert) · Kaveri IT Cell Lead (IT Security / Engineering)",
        size=14.5,
    )
    add_rect(slide, MARGIN, Inches(6.14), CONTENT_W, Inches(0.66), fill=BAND, line=RULE)
    add_text(
        slide, MARGIN + Inches(0.24), Inches(6.30), CONTENT_W - Inches(0.48), Inches(0.4),
        "Source of record: BRD_User_Management_v4.18.docx  ·  Editable diagram sources: ProcessDiagrams/User_Management/*.drawio",
        size=11.5, color=MUTED,
    )

    return deck.prs


def main() -> None:
    if not DOCX.exists():
        raise FileNotFoundError(DOCX)
    workdir = Path(tempfile.mkdtemp(prefix="um_ppt_"))
    try:
        figures = extract_figures(workdir)
        prs = build(figures)
        target = DST
        try:
            prs.save(str(target))
        except PermissionError:
            target = DST.with_name(DST.stem + "_unlocked" + DST.suffix)
            prs.save(str(target))
            print("ORIGINAL LOCKED (open in PowerPoint) — saved instead as:")
        print(f"{target} ({target.stat().st_size} bytes, {len(prs.slides)} slides)")
        claude_dir = BASE.parent.parent / "Claude"
        if claude_dir.is_dir():
            try:
                shutil.copy2(target, claude_dir / target.name)
                print(f"Mirrored: {claude_dir / target.name}")
            except Exception as exc:
                print(f"Claude mirror skipped: {exc}")
    finally:
        shutil.rmtree(workdir, ignore_errors=True)


if __name__ == "__main__":
    main()
